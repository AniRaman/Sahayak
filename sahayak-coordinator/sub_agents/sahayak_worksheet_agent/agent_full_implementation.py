"""
Full implementation of educational materials generation functions.
This file contains all the helper functions used by the mcp_server.
"""

import os
import json
import logging
from typing import Dict, List, Any, Optional
from datetime import datetime, timedelta
import uuid
import base64
from io import BytesIO

# File generation libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_LEFT, TA_CENTER

# Document processing libraries
import pdfplumber
from docx import Document
import pytesseract
from PIL import Image

# Google Cloud and ADK imports
from google.cloud import firestore
from google.cloud import storage
from google.cloud import vision
from google.oauth2 import service_account
from google.adk.tools.tool_context import ToolContext
from google.adk.agents.callback_context import CallbackContext
import google.generativeai as genai
import google.genai.types as types

# Pipeline import - using simple pipeline for now
try:
    from .simple_pipeline import process_educational_content_simple as process_educational_content
except ImportError:
    try:
        from simple_pipeline import process_educational_content_simple as process_educational_content
    except ImportError:
        def process_educational_content(content, context=None):
            return {
                "status": "error", 
                "error_message": "Simple pipeline not available - using ADK processing only"
            }

# Environment setup
from dotenv import load_dotenv
load_dotenv(override=True)

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)

def _get_authenticated_storage_client():
    """Get authenticated Google Cloud Storage client"""
    project_id = os.getenv("GOOGLE_CLOUD_PROJECT_ID")
    
    if not project_id:
        raise ValueError("GOOGLE_CLOUD_PROJECT_ID environment variable not set")
    
    # Try Sahayak-specific service account first
    sahayak_service_account_path = os.getenv("SAHAYAK_SERVICE_ACCOUNT_CREDENTIALS")
    if sahayak_service_account_path and os.path.exists(sahayak_service_account_path):
        print(f"[AUTH] Using Sahayak service account file: {sahayak_service_account_path}")
        credentials = service_account.Credentials.from_service_account_file(sahayak_service_account_path)
        return storage.Client(project=project_id, credentials=credentials)
    
    # Fallback to default credentials
    return storage.Client(project=project_id)


async def generate_educational_materials(
    content: str = "",
    subject: str = "",
    grade_level: str = "",
    curriculum: str = "",
    output_format: str = "both",  # "slides", "worksheets", "both"
    save_to_storage: bool = True
) -> Dict[str, Any]:
    """
    Main function to generate educational materials from content.
    """
    
    print(f"[WORKSHEET_AGENT] Processing educational content for {subject} {grade_level}")
    print("content :",content)
    try:
        # Validate final content
        if not content.strip():
            return {
                "status": "error", 
                "error_message": "No content provided. Please provide text content.",
                "timestamp": datetime.now().isoformat()
            }
        
        # Step 2: Process content through sequential pipeline
        additional_context = {
            "subject": subject,
            "grade_level": grade_level,
            "curriculum": curriculum
        }
        
        pipeline_result = process_educational_content(content, additional_context)
        print("pipeline_result :",pipeline_result)
        print(f"[PIPELINE_DEBUG] Pipeline result status: {pipeline_result.get('status')}")
        print(f"[PIPELINE_DEBUG] Pipeline result keys: {list(pipeline_result.keys())}")
        
        if pipeline_result.get("status") != "success":
            return pipeline_result
        
        # Extract pipeline outputs
        content_analysis = pipeline_result.get("content_analysis", {})
        slide_structure = pipeline_result.get("slide_structure", {})
        worksheets_data = pipeline_result.get("worksheets", {})
        
        print(f"[PIPELINE_DEBUG] Content analysis keys: {list(content_analysis.keys()) if content_analysis else 'None'}")
        print(f"[PIPELINE_DEBUG] Slide structure keys: {list(slide_structure.keys()) if slide_structure else 'None'}")
        print(f"[PIPELINE_DEBUG] Worksheets data keys: {list(worksheets_data.keys()) if worksheets_data else 'None'}")
        
        # Parse raw_response fields if they exist (pipeline returns JSON strings)
        import json
        
        if content_analysis and "raw_response" in content_analysis:
            try:
                print(f"[PIPELINE_DEBUG] Parsing content_analysis raw_response...")
                raw_response = content_analysis["raw_response"]
                # Extract JSON from markdown code block if present
                if "```json" in raw_response:
                    json_start = raw_response.find("```json") + 7
                    json_end = raw_response.find("```", json_start)
                    raw_response = raw_response[json_start:json_end].strip()
                
                parsed_content = json.loads(raw_response)
                content_analysis = parsed_content
                print(f"[PIPELINE_DEBUG] Parsed content_analysis: {list(content_analysis.keys())}")
            except Exception as e:
                print(f"[PIPELINE_DEBUG] Failed to parse content_analysis raw_response: {e}")
                content_analysis = {}
        
        if slide_structure and "raw_response" in slide_structure:
            try:
                print(f"[PIPELINE_DEBUG] Parsing slide_structure raw_response...")
                raw_response = slide_structure["raw_response"]
                # Extract JSON from markdown code block if present
                if "```json" in raw_response:
                    json_start = raw_response.find("```json") + 7
                    json_end = raw_response.find("```", json_start)
                    raw_response = raw_response[json_start:json_end].strip()
                
                parsed_slides = json.loads(raw_response)
                slide_structure = parsed_slides
                print(f"[PIPELINE_DEBUG] Parsed slide_structure: {list(slide_structure.keys())}")
            except Exception as e:
                print(f"[PIPELINE_DEBUG] Failed to parse slide_structure raw_response: {e}")
                slide_structure = {}
        
        if worksheets_data and "raw_response" in worksheets_data:
            try:
                print(f"[PIPELINE_DEBUG] Parsing worksheets_data raw_response...")
                raw_response = worksheets_data["raw_response"]
                # Extract JSON from markdown code block if present
                if "```json" in raw_response:
                    json_start = raw_response.find("```json") + 7
                    json_end = raw_response.find("```", json_start)
                    raw_response = raw_response[json_start:json_end].strip()
                
                parsed_worksheets = json.loads(raw_response)
                worksheets_data = parsed_worksheets
                print(f"[PIPELINE_DEBUG] Parsed worksheets_data: {list(worksheets_data.keys())}")
            except Exception as e:
                print(f"[PIPELINE_DEBUG] Failed to parse worksheets_data raw_response: {e}")
                worksheets_data = {}
        
        # If pipeline didn't return proper data structures, create fallback ones
        if not content_analysis or not isinstance(content_analysis, dict):
            content_analysis = {
                "subject": subject or "General",
                "topic": "Educational Content",
                "grade_level": grade_level or "General",
                "difficulty": "intermediate"
            }
            print(f"[PIPELINE_DEBUG] Using fallback content_analysis")
        
        if not slide_structure or not slide_structure.get("slides"):
            slide_structure = {
                "slides": [
                    {
                        "slide_type": "title",
                        "title": content_analysis.get("topic", "Educational Content"),
                        "subtitle": f"Subject: {content_analysis.get('subject', 'General')}"
                    },
                    {
                        "slide_type": "content",
                        "title": "Key Points",
                        "content": content[:500].split('\n')[:5]  # First 5 lines or 500 chars
                    }
                ]
            }
            print(f"[PIPELINE_DEBUG] Using fallback slide_structure with {len(slide_structure['slides'])} slides")
        
        if not worksheets_data or not worksheets_data.get("worksheets"):
            # Create fallback worksheet data
            worksheets_data = {
                "worksheets": {
                    "easy": {
                        "instructions": "Answer the following questions based on the content.",
                        "estimated_time": "20 minutes",
                        "questions": [
                            {
                                "question_number": 1,
                                "question": "What is the main topic of this content?",
                                "question_type": "short_answer",
                                "points": 5
                            },
                            {
                                "question_number": 2,
                                "question": "List three key points from the content.",
                                "question_type": "short_answer",
                                "points": 10
                            }
                        ]
                    },
                    "medium": {
                        "instructions": "Answer the following questions with detailed explanations.",
                        "estimated_time": "30 minutes",
                        "questions": [
                            {
                                "question_number": 1,
                                "question": "Explain the main concepts discussed in the content.",
                                "question_type": "short_answer",
                                "points": 10
                            },
                            {
                                "question_number": 2,
                                "question": "How can you apply these concepts in real life?",
                                "question_type": "short_answer",
                                "points": 15
                            }
                        ]
                    },
                    "hard": {
                        "instructions": "Provide comprehensive answers with examples and analysis.",
                        "estimated_time": "45 minutes",
                        "questions": [
                            {
                                "question_number": 1,
                                "question": "Analyze and critically evaluate the content presented.",
                                "question_type": "short_answer",
                                "points": 20
                            },
                            {
                                "question_number": 2,
                                "question": "Create your own example demonstrating the concepts.",
                                "question_type": "short_answer",
                                "points": 20
                            }
                        ]
                    }
                }
            }
            print(f"[PIPELINE_DEBUG] Using fallback worksheets_data with {len(worksheets_data['worksheets'])} difficulty levels")
        
        # Step 2: Generate files based on output_format
        generated_files = {}
        
        if output_format in ["slides", "both"]:
            print(f"[FILE_GEN] Generating presentation slides...")
            print(f"[FILE_GEN] Slide structure has {len(slide_structure.get('slides', []))} slides")
            slides_result = _generate_presentation_file(slide_structure, content_analysis)
            print(f"[FILE_GEN] Slides result status: {slides_result.get('status')}")
            if slides_result.get("status") == "success":
                generated_files["presentation"] = slides_result
                print(f"[FILE_GEN] Added presentation to generated files")
            else:
                print(f"[FILE_GEN] Slides generation failed: {slides_result.get('error_message')}")
        
        if output_format in ["worksheets", "both"]:
            print(f"[FILE_GEN] Generating differentiated worksheets...")
            print(f"[FILE_GEN] Worksheets data has {len(worksheets_data.get('worksheets', {}))} difficulty levels")
            worksheets_result = _generate_worksheet_files(worksheets_data, content_analysis)
            print(f"[FILE_GEN] Worksheets result status: {worksheets_result.get('status')}")
            print(f"[FILE_GEN] Worksheets result files: {list(worksheets_result.get('files', {}).keys())}")
            if worksheets_result.get("status") == "success":
                generated_files.update(worksheets_result.get("files", {}))
                print(f"[FILE_GEN] Added {len(worksheets_result.get('files', {}))} worksheets to generated files")
            else:
                print(f"[FILE_GEN] Worksheets generation failed: {worksheets_result.get('error_message')}")
        
        print(f"[FILE_GEN] Total generated files: {len(generated_files)} - {list(generated_files.keys())}")
        
        # Step 3: Save to Google Cloud Storage and Firestore (if requested)
        final_result = {
            "status": "success",
            "content_analysis": content_analysis,
            "generated_materials": {},
            "download_links": {},
            "metadata": {
                "generation_timestamp": datetime.now().isoformat(),
                "subject": content_analysis.get("subject", subject),
                "topic": content_analysis.get("topic", ""),
                "grade_level": content_analysis.get("grade_level", grade_level),
                "difficulty": content_analysis.get("difficulty", "intermediate"),
                "curriculum": curriculum,
                "files_generated": list(generated_files.keys())
            }
        }
        
        if save_to_storage and generated_files:
            print(f"[STORAGE] Saving {len(generated_files)} files to Google Cloud Storage...")
            storage_result = _save_files_to_storage(generated_files, content_analysis)
            
            if storage_result.get("status") == "success":
                final_result["download_links"] = storage_result.get("download_links", {})
                final_result["storage_metadata"] = storage_result.get("storage_metadata", {})
                
                # Save metadata to Firestore for easy retrieval
                firestore_result = _save_metadata_to_firestore(final_result)
                if firestore_result.get("status") == "success":
                    final_result["firestore_document_id"] = firestore_result.get("document_id")
        else:
            # Return file content without saving
            final_result["generated_materials"] = generated_files
        
        print(f"[WORKSHEET_SUCCESS] Generated {len(generated_files)} educational materials")
        return final_result
        
    except Exception as e:
        logger.error(f"Educational material generation failed: {e}")
        return {
            "status": "error",
            "error_message": f"Material generation failed: {str(e)}",
            "timestamp": datetime.now().isoformat()
        }


def _generate_presentation_file(slide_structure: Dict, content_analysis: Dict) -> Dict[str, Any]:
    """Generate PowerPoint presentation file from slide structure"""
    
    try:
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 widescreen)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        slides_data = slide_structure.get("slides", [])
        
        for slide_info in slides_data:
            slide_type = slide_info.get("slide_type", "content")
            
            if slide_type == "title":
                slide = _create_title_slide(prs, slide_info)
            elif slide_type == "objectives":
                slide = _create_bullet_slide(prs, slide_info)
            else:
                slide = _create_content_slide(prs, slide_info)
        
        # Save to temporary file
        temp_filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        temp_path = os.path.join(os.getcwd(), temp_filename)
        
        prs.save(temp_path)
        
        # Read file content for storage
        with open(temp_path, 'rb') as f:
            file_content = f.read()
        
        # Clean up temp file
        os.remove(temp_path)
        
        presentation_info = {
            "status": "success",
            "filename": f"{content_analysis.get('topic', 'Educational_Content')}_Slides.pptx",
            "content": file_content,
            "size_bytes": len(file_content),
            "slides_count": len(slides_data),
            "file_type": "presentation",
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        }
        
        return presentation_info
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Presentation generation failed: {str(e)}"
        }


def _create_title_slide(prs: Presentation, slide_info: Dict):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = slide_info.get("title", "Educational Presentation")
    subtitle.text = slide_info.get("subtitle", "")
    
    return slide


def _create_bullet_slide(prs: Presentation, slide_info: Dict):
    """Create slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = slide_info.get("title", "")
    
    # Add bullet points
    text_frame = content.text_frame
    text_frame.clear()
    
    for i, point in enumerate(slide_info.get("content", [])):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
    
    return slide


def _create_content_slide(prs: Presentation, slide_info: Dict):
    """Create general content slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = slide_info.get("title", "")
    
    # Add content
    if slide_info.get("content"):
        content = slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        
        for i, point in enumerate(slide_info.get("content", [])):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
    
    return slide


def _generate_worksheet_files(worksheets_data: Dict, content_analysis: Dict) -> Dict[str, Any]:
    """Generate PDF worksheet files for each difficulty level"""
    
    try:
        worksheets = worksheets_data.get("worksheets", {})
        generated_files = {}
        
        for difficulty in ["easy", "medium", "hard"]:
            if difficulty in worksheets:
                worksheet_info = worksheets[difficulty]
                
                # Generate PDF for this difficulty level
                pdf_result = _create_worksheet_pdf(worksheet_info, difficulty, content_analysis)
                
                if pdf_result.get("status") == "success":
                    file_key = f"worksheet_{difficulty}"
                    generated_files[file_key] = pdf_result
        
        return {
            "status": "success",
            "files": generated_files,
            "worksheets_count": len(generated_files)
        }
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Worksheet generation failed: {str(e)}"
        }


def _create_worksheet_pdf(worksheet_info: Dict, difficulty: str, content_analysis: Dict) -> Dict[str, Any]:
    """Create PDF file for a single worksheet"""
    
    try:
        # Create temporary PDF file
        temp_filename = f"worksheet_{difficulty}_{uuid.uuid4().hex[:8]}.pdf"
        temp_path = os.path.join(os.getcwd(), temp_filename)
        
        # Create PDF document
        doc = SimpleDocTemplate(temp_path, pagesize=A4)
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            alignment=TA_CENTER,
            spaceBefore=0,
            spaceAfter=12
        )
        
        # Build PDF content
        story = []
        
        # Title
        topic = content_analysis.get("topic", "Educational Worksheet")
        title = f"{topic} - {difficulty.title()} Level"
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 12))
        
        # Instructions
        instructions = worksheet_info.get("instructions", "Complete the following exercises.")
        story.append(Paragraph(f"<b>Instructions:</b> {instructions}", styles['Normal']))
        
        # Time estimate
        time_estimate = worksheet_info.get("estimated_time", "30 minutes")
        story.append(Paragraph(f"<b>Estimated Time:</b> {time_estimate}", styles['Normal']))
        story.append(Spacer(1, 12))
        
        # Questions
        questions = worksheet_info.get("questions", [])
        
        for i, question in enumerate(questions):
            # Question number and text
            q_num = question.get("question_number", i + 1)
            q_text = question.get("question", "")
            points = question.get("points", 1)
            
            story.append(Paragraph(f"<b>Question {q_num}:</b> {q_text} <i>({points} points)</i>", styles['Normal']))
            
            # Handle different question types
            q_type = question.get("question_type", "short_answer")
            
            if q_type == "multiple_choice":
                options = question.get("options", [])
                for option in options:
                    story.append(Paragraph(f"   {option}", styles['Normal']))
            
            # Add space for answer
            story.append(Spacer(1, 18))
            story.append(Paragraph("Answer: _" + "_" * 50, styles['Normal']))
            story.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(story)
        
        # Read file content
        with open(temp_path, 'rb') as f:
            file_content = f.read()
        
        # Clean up temp file
        os.remove(temp_path)
        
        worksheet_result = {
            "status": "success",
            "filename": f"{topic.replace(' ', '_')}_{difficulty}_worksheet.pdf",
            "content": file_content,
            "size_bytes": len(file_content),
            "difficulty": difficulty,
            "questions_count": len(questions),
            "file_type": "worksheet",
            "content_type": "application/pdf"
        }
        
        return worksheet_result
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"PDF generation failed: {str(e)}"
        }


def _save_files_to_storage(generated_files: Dict, content_analysis: Dict) -> Dict[str, Any]:
    """Save generated files to Google Cloud Storage"""
    
    try:
        storage_client = _get_authenticated_storage_client()
        bucket_name = os.getenv("GOOGLE_CLOUD_STORAGE_BUCKET", "agentic_storage")
        # Clean bucket name to ensure no escaping issues
        bucket_name = bucket_name.replace("\\_", "_").replace("\\", "")
        print(f"[STORAGE_DEBUG] Final bucket name: '{bucket_name}'")
        bucket = storage_client.bucket(bucket_name)
        
        download_links = {}
        storage_metadata = {}
        
        # Create folder structure: worksheets/subject/topic/
        subject = content_analysis.get("subject", "General").replace(" ", "_")
        topic = content_analysis.get("topic", "Educational_Content").replace(" ", "_")
        folder_path = f"generated_materials/{subject}/{topic}"
        
        for file_key, file_info in generated_files.items():
            if file_info.get("status") == "success":
                filename = file_info.get("filename", f"{file_key}.pdf")
                file_content = file_info.get("content")
                content_type = file_info.get("content_type", "application/pdf")
                
                # Upload to storage
                blob_path = f"{folder_path}/{filename}"
                blob = bucket.blob(blob_path)
                
                blob.upload_from_string(file_content, content_type=content_type)
                
                # Generate signed URL for download
                print(f"[STORAGE_DEBUG] Bucket name: '{bucket_name}'")
                print(f"[STORAGE_DEBUG] Blob path: '{blob_path}'")
                try:
                    download_url = blob.generate_signed_url(
                        version="v4",
                        expiration=timedelta(hours=24),
                        method="GET"
                    )
                    print(f"[STORAGE_SUCCESS] Generated signed URL for {filename}")
                    print(f"[STORAGE_DEBUG] Signed URL: {download_url}")
                except Exception as url_error:
                    print(f"[STORAGE_ERROR] Failed to generate signed URL for {filename}: {str(url_error)}")
                    # Use public URL as fallback (if bucket allows public access)
                    # Make sure bucket name is not escaped
                    clean_bucket_name = bucket_name.replace("\\_", "_")
                    download_url = f"https://storage.googleapis.com/{clean_bucket_name}/{blob_path}"
                    print(f"[STORAGE_DEBUG] Fallback URL: {download_url}")
                
                download_links[file_key] = {
                    "filename": filename,
                    "download_url": download_url,
                    "file_type": file_info.get("file_type"),
                    "size_bytes": file_info.get("size_bytes"),
                    "storage_path": blob_path
                }
                
                print(f"[STORAGE_SUCCESS] Saved {filename} to {blob_path}")
        
        storage_metadata = {
            "bucket": bucket_name,
            "folder_path": folder_path,
            "files_uploaded": len(download_links),
            "upload_timestamp": datetime.now().isoformat()
        }
        
        return {
            "status": "success",
            "download_links": download_links,
            "storage_metadata": storage_metadata
        }
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Storage upload failed: {str(e)}"
        }


def _save_metadata_to_firestore(generation_result: Dict) -> Dict[str, Any]:
    """Save generation metadata to Firestore for easy retrieval"""
    
    try:
        project_id = os.getenv("GOOGLE_CLOUD_PROJECT_ID")
        database_id = os.getenv("FIRESTORE_DATABASE_ID", "agenticdb")
        
        if not project_id:
            return {"status": "error", "error_message": "Project ID not configured"}
        
        db = firestore.Client(project=project_id, database=database_id)
        
        # Prepare metadata document
        metadata = generation_result.get("metadata", {})
        content_analysis = generation_result.get("content_analysis", {})
        download_links = generation_result.get("download_links", {})
        
        doc_data = {
            "title": f"{content_analysis.get('topic', 'Educational Materials')} - Generated Materials",
            "subject": content_analysis.get("subject", "General"),
            "topic": content_analysis.get("topic", ""),
            "grade_level": content_analysis.get("grade_level", ""),
            "difficulty": content_analysis.get("difficulty", "intermediate"),
            "curriculum": metadata.get("curriculum", ""),
            "material_type": "generated_worksheet_package",
            "files": [
                {
                    "type": link_info.get("file_type"),
                    "filename": link_info.get("filename"),
                    "download_url": link_info.get("download_url"),
                    "size_bytes": link_info.get("size_bytes", 0)
                }
                for link_info in download_links.values()
            ],
            "generation_metadata": {
                "generated_at": metadata.get("generation_timestamp"),
                "content_length": len(str(content_analysis)),
                "files_generated": metadata.get("files_generated", [])
            },
            "created_date": firestore.SERVER_TIMESTAMP,
            "access_count": 0,
            "rating": 5.0  # Default high rating for generated content
        }
        
        # Save to Firestore
        doc_ref = db.collection("generated_materials").add(doc_data)
        document_id = doc_ref[1].id
        
        print(f"[FIRESTORE_SUCCESS] Saved metadata with ID: {document_id}")
        
        return {
            "status": "success",
            "document_id": document_id,
            "collection": "generated_materials"
        }
        
    except Exception as e:
        return {
            "status": "error",
            "error_message": f"Firestore save failed: {str(e)}"
        }